# indexer.py
# Async, cancellable indexing engine for the MCP server.
# Extracted from index_docs.py logic, designed to run as a background task.

import asyncio
import hashlib
import json
import logging
import os
import time
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from langchain_core.documents import Document

import config
import logging_config
import settings as settings_module


STATE_FILE = Path(__file__).parent / "indexing_state.json"
ERROR_FILES_FILE = Path(__file__).parent / "indexing_errors.json"
CONFIG_SNAPSHOT_FILE = Path(__file__).parent / "index_config_snapshot.json"


class EmbeddingCircuitBreaker(Exception):
    """Raised when too many consecutive embedding failures for one file; skip file and continue."""


def _load_state() -> Dict[str, Dict[str, Any]]:
    if STATE_FILE.exists():
        try:
            with open(STATE_FILE, "r") as f:
                return json.load(f)
        except Exception as e:
            logging.error(f"Failed to load state file: {e}")
    return {}


def _save_state(state: Dict[str, Dict[str, Any]]):
    with open(STATE_FILE, "w") as f:
        json.dump(state, f, indent=2)


def _load_error_files() -> List[str]:
    """Load persisted file paths that had errors in the last indexing run."""
    if ERROR_FILES_FILE.exists():
        try:
            with open(ERROR_FILES_FILE, "r") as f:
                data = json.load(f)
                return data.get("files", [])
        except Exception as e:
            logging.warning(f"Failed to load error files: {e}")
    return []


def _save_error_files(files: List[str]):
    """Persist file paths that had errors (survives restarts)."""
    with open(ERROR_FILES_FILE, "w") as f:
        json.dump({"files": files, "timestamp": time.time()}, f, indent=2)


def _save_index_config_snapshot(directories: List[str], extensions: List[str], exclusions: List[str]):
    """Save indexing-relevant config when indexing completes (for dirty-flag comparison)."""
    try:
        with open(CONFIG_SNAPSHOT_FILE, "w") as f:
            json.dump({
                "directories": sorted(directories) if directories else [],
                "extensions": sorted(extensions) if extensions else [],
                "exclusions": sorted(exclusions) if exclusions else [],
            }, f, indent=2)
    except Exception as e:
        logging.warning(f"Failed to save index config snapshot: {e}")


def _load_index_config_snapshot() -> Optional[Dict[str, List[str]]]:
    """Load the config snapshot from last successful index."""
    if not CONFIG_SNAPSHOT_FILE.exists():
        return None
    try:
        with open(CONFIG_SNAPSHOT_FILE, "r") as f:
            data = json.load(f)
            return {
                "directories": data.get("directories", []),
                "extensions": data.get("extensions", []),
                "exclusions": data.get("exclusions", []),
            }
    except Exception as e:
        logging.warning(f"Failed to load index config snapshot: {e}")
        return None


def _read_file_content(file_path: Path) -> str:
    """Read text content from a file. Uses lightweight, CPU-only extractors for office docs."""
    suffix = file_path.suffix.lower()

    # Fast path: plain text and code
    if suffix in [".md", ".txt", ".py", ".js", ".json", ".sh", ".css", ".sql", ".yaml", ".yml"]:
        encodings = ["utf-8", "latin-1", "cp1252"]
        for encoding in encodings:
            try:
                with open(file_path, "r", encoding=encoding) as f:
                    return f.read()
            except UnicodeDecodeError:
                continue
            except Exception as e:
                logging.warning(f"Error reading {file_path}: {e}")
                return ""
        return ""

    # PDF: PyMuPDF (fitz) — fast, preserves order and structure
    if suffix == ".pdf":
        try:
            import fitz
            with fitz.open(file_path) as doc:
                parts = []
                for page in doc:
                    text = page.get_text(sort=True)
                    if text:
                        parts.append(text)
                return "\n\n".join(parts) if parts else ""
        except Exception as e:
            logging.warning(f"PyMuPDF failed for {file_path}: {type(e).__name__}: {e}")
            return ""

    # DOCX: python-docx — paragraphs, headings, lists, tables
    if suffix == ".docx":
        try:
            from docx import Document
            doc = Document(file_path)
            parts = []
            for para in doc.paragraphs:
                if para.text.strip():
                    parts.append(para.text)
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    if row_text:
                        parts.append(row_text)
            return "\n\n".join(parts) if parts else ""
        except Exception as e:
            logging.warning(f"python-docx failed for {file_path}: {type(e).__name__}: {e}")
            return ""

    # PPTX: python-pptx — slides, notes, shapes, tables, grouped shapes (recursive)
    if suffix == ".pptx":
        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE

            def _extract_shape_text(shape, parts: list) -> None:
                """Extract text from a shape or its grouped children."""
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    for child in shape.shapes:
                        _extract_shape_text(child, parts)
                elif shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            parts.append(para.text.strip())
                elif shape.has_table:
                    for row in shape.table.rows:
                        row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                        if row_text:
                            parts.append(row_text)

            prs = Presentation(file_path)
            parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    _extract_shape_text(shape, parts)
                try:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        parts.append(notes)
                except AttributeError:
                    pass
            return "\n\n".join(parts) if parts else ""
        except Exception as e:
            logging.warning(f"python-pptx failed for {file_path}: {type(e).__name__}: {e}")
            return ""

    # XLSX: openpyxl — sheet text (data_only=False avoids slow formula evaluation)
    if suffix == ".xlsx":
        try:
            from openpyxl import load_workbook
            wb = load_workbook(file_path, read_only=True, data_only=False, keep_vba=False)
            parts = []
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    row_text = " | ".join(str(c) for c in row if c is not None and str(c).strip())
                    if row_text.strip():
                        parts.append(row_text)
            wb.close()
            return "\n\n".join(parts) if parts else ""
        except Exception as e:
            logging.warning(f"openpyxl failed for {file_path}: {type(e).__name__}: {e}")
            return ""

    # HTML: html2text — convert to plain text
    if suffix == ".html" or suffix == ".htm":
        try:
            import html2text
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                html = f.read()
            h = html2text.HTML2Text()
            h.ignore_links = False
            h.ignore_images = True
            h.body_width = 0
            return h.handle(html)
        except Exception as e:
            logging.warning(f"html2text failed for {file_path}: {type(e).__name__}: {e}")
            return ""

    return ""


def _scan_files(directories: List[str], extensions: List[str], exclusions: List[str]) -> List[Tuple[Path, float]]:
    """Recursively find all files with matching extensions, respecting exclusions.
    Returns list of (Path, mtime) tuples.
    """
    files: List[Tuple[Path, float]] = []
    
    # Pre-process exclusions for faster lookup
    # sets are faster than lists
    exclude_names = set(exclusions)
    
    # Create a thread pool for parallel scanning
    # Limit max workers to avoid excessive threads, but enough to cover I/O latency
    import concurrent.futures
    
    # Check for Windows Optimization
    user_settings = settings_module.load_settings()
    if user_settings.get("optimize_for_windows", False):
         try:
             import windows.bridge as win_bridge
             logging.info("Attempting Windows Native Scan...")
             bridge_files = win_bridge.scan_windows(directories, extensions, list(exclude_names))
             if bridge_files is not None:
                 logging.info(f"Windows Native Scan successful: {len(bridge_files)} files.")
                 return bridge_files
             else:
                 logging.warning("Windows Native Scan returned None, falling back to parallel scan.")
         except Exception as e:
             logging.error(f"Failed to use Windows Native Scan: {e}")

    max_workers = min(32, (os.cpu_count() or 1) * 4) 
    
    with concurrent.futures.ThreadPoolExecutor(max_workers=max_workers) as executor:
        # Map of future -> directory path
        futures = {}
        
        # Helper to submit a directory scan task
        def submit_scan(dir_path: Path):
            return executor.submit(_scan_directory_task, dir_path, exclude_names, extensions)

        # Initial submissions
        for directory in directories:
            path = Path(directory)
            if not path.exists():
                logging.warning(f"Directory not found: {directory}")
                continue
            futures[submit_scan(path)] = path
            
        while futures:
            # Wait for any future to complete
            done, _ = concurrent.futures.wait(futures, return_when=concurrent.futures.FIRST_COMPLETED)
            
            for future in done:
                path = futures.pop(future)
                try:
                    result_files, result_subdirs = future.result()
                    files.extend(result_files)
                    
                    # Submit subdirectories
                    for subdir in result_subdirs:
                         futures[submit_scan(subdir)] = subdir
                         
                except Exception as e:
                    logging.warning(f"Error scanning {path}: {e}")

    return files

def _scan_directory_task(path: Path, exclude_names: set, extensions: List[str]) -> Tuple[List[Tuple[Path, float]], List[Path]]:
    """
    Scans a single directory. Returns (files, subdirectories).
    Executed in a thread pool.
    """
    local_files = []
    subdirs = []
    try:
        with os.scandir(path) as it:
            # Consume iterator
            entries = list(it)
            
            for entry in entries:
                if entry.name.startswith("."):
                    continue
                
                if entry.is_dir():
                    if entry.name not in exclude_names:
                        subdirs.append(Path(entry.path))
                    continue

                # File processing
                if (entry.name.endswith(".log") or 
                    entry.name.endswith(".log.txt") or 
                    "_logs" in entry.name or 
                    entry.name.startswith("log-")):
                    continue

                _, ext = os.path.splitext(entry.name)
                if ext.lower() not in extensions:
                        continue

                try:
                    stat = entry.stat()
                    if stat.st_size > 5 * 1024 * 1024:
                            continue
                            
                    local_files.append((Path(entry.path), stat.st_mtime))
                except OSError:
                    pass
                    
    except OSError as e:
        # Don't log here to avoid thread contention on logging lock? 
        # Or just let the main thread handle exception?
        # Re-raising allows main loop to log associated with path
        raise e
        
    return local_files, subdirs


def unload_llm():
    """Explicitly unload the chat model to free VRAM for embeddings."""
    try:
        import requests
        import config
        # Payload keep_alive: 0 triggers unload
        model = config.LLM_MODEL
        logging.info(f"Unloading LLM {model} to free VRAM...")
        try:
             requests.post(f"{config.OLLAMA_BASE_URL}/api/generate", 
                           json={"model": model, "keep_alive": 0}, timeout=2)
        except Exception:
             pass
    except Exception as e:
        logging.warning(f"Failed to unload LLM: {e}")


def _quiesce_services_for_indexing():
    """Stop Open WebUI and MCP so Ollama only runs embeddings during indexing."""
    import subprocess
    script_dir = Path(__file__).parent.resolve()
    try:
        logging.info("Quiescing services for indexing: stopping Open WebUI and MCP...")
        # Stop MCP server (prevents chat/RAG requests)
        subprocess.run(
            ["/bin/bash", str(script_dir / "run_mcp_server.sh"), "stop"],
            cwd=str(script_dir), timeout=5, capture_output=True,
        )
        # Stop Open WebUI container (prevents chat UI from loading LLM)
        cmd = subprocess.run(
            ["sh", "-c", "command -v podman >/dev/null && podman stop open-webui 2>/dev/null || docker stop open-webui 2>/dev/null || true"],
            cwd=str(script_dir), timeout=5, capture_output=True,
        )
        logging.info("Services quiesced. Ollama will run embeddings only.")
        time.sleep(2)  # Let WebUI/MCP shut down before embedding
    except Exception as e:
        logging.warning(f"Failed to quiesce services: {e}")


def _restore_services_after_indexing():
    """Warm LLM and restart Open WebUI + MCP so chat is responsive after indexing."""
    import subprocess
    import threading

    def _run():
        try:
            import requests
            import config
            # 1. Warm the LLM (load into memory for chat)
            model = config.LLM_MODEL
            logging.info(f"Restoring services: warming LLM {model}...")
            try:
                requests.post(
                    f"{config.OLLAMA_BASE_URL}/api/generate",
                    json={"model": model, "prompt": "hi", "stream": False, "keep_alive": "30m"},
                    timeout=60,
                )
                logging.info("LLM warmed.")
            except Exception as e:
                logging.warning(f"Failed to warm LLM: {e}")

            # 2. Restart Open WebUI and MCP (may have become unresponsive under load)
            script_dir = Path(__file__).parent.resolve()
            try:
                logging.info("Restoring services: restarting MCP server...")
                subprocess.run(
                    ["/bin/bash", str(script_dir / "run_mcp_server.sh"), "stop"],
                    cwd=str(script_dir), timeout=5, capture_output=True,
                )
                subprocess.run(
                    ["/bin/bash", str(script_dir / "run_mcp_server.sh"), "start"],
                    cwd=str(script_dir), timeout=10, capture_output=True,
                )
                logging.info("MCP server restarted.")
            except Exception as e:
                logging.warning(f"Failed to restart MCP: {e}")

            try:
                logging.info("Restoring services: restarting Open WebUI...")
                subprocess.run(
                    ["/bin/bash", str(script_dir / "run_webui.sh")],
                    cwd=str(script_dir), timeout=30, capture_output=True,
                )
                logging.info("Open WebUI restarted.")
            except Exception as e:
                logging.warning(f"Failed to restart Open WebUI: {e}")

        except Exception as e:
            logging.warning(f"Post-indexing restore failed: {e}")

    t = threading.Thread(target=_run, daemon=True)
    t.start()

class IndexingJob:
    """Manages a single indexing run with progress tracking and cancellation."""

    def __init__(self):
        self.status: str = "idle"  # idle | preparing | scanning | indexing | complete | cancelled | error
        self.total_files: int = 0
        self.files_to_process: int = 0
        self.processed_files: int = 0
        self.current_file: str = ""
        self.errors: List[str] = []
        self.error_files: List[str] = []  # Full paths of files that had errors (for re-index errors only)
        self.files_deleted: int = 0
        self._cancelled: bool = False
        self.detail_status: str = ""  # Granular status message (e.g. "Loading models...")
        self._start_time: Optional[float] = None
        self.parse_time_seconds: float = 0.0  # Cumulative time in file parsing + chunking
        self.embed_time_seconds: float = 0.0   # Cumulative time in embedding + Qdrant

    def cancel(self):
        """Request cancellation of the current job."""
        if self.status in ("preparing", "scanning", "indexing"):
            self._cancelled = True
            self.status = "cancelled"

    def get_status(self) -> dict:
        """Return current job status as a dict."""
        result = {
            "status": self.status,
            "total_files_found": self.total_files,
            "files_to_process": self.files_to_process,
            "processed_files": self.processed_files,
            "current_file": self.current_file,
            "errors": len(self.errors),
            "error_details": self.errors[-20:],  # Last 20 errors for UI
            "error_files_count": len(self.error_files) or len(_load_error_files()),  # For re-index errors button
            "files_deleted": self.files_deleted,
            "detail_status": self.detail_status,
        }
        if self._start_time and self.status == "indexing" and self.files_to_process > 0:
            elapsed = time.time() - self._start_time
            if self.processed_files > 0:
                rate = elapsed / self.processed_files
                remaining = (self.files_to_process - self.processed_files) * rate
                result["elapsed_seconds"] = round(elapsed)
                result["eta_seconds"] = round(remaining)
                result["percent"] = round(100 * self.processed_files / self.files_to_process)
        result["parse_time_seconds"] = round(self.parse_time_seconds)
        result["embed_time_seconds"] = round(self.embed_time_seconds)
        return result

    def get_status_text(self) -> str:
        """Return a human-readable progress string."""
        s = self.get_status()
        if s["status"] == "idle":
            return "No indexing job running."
        if s["status"] == "scanning":
            return "Scanning directories for files..."
        if s["status"] == "cancelled":
            return f"Indexing cancelled. Processed {s['processed_files']}/{s['files_to_process']} files."
        if s["status"] == "error":
            return f"Indexing failed: {self.errors[-1] if self.errors else 'unknown error'}"
        if s["status"] == "complete":
            return (
                f"Indexing complete. {s['processed_files']} files indexed, "
                f"{s['files_deleted']} deleted from index. "
                f"{s['errors']} errors."
            )
        # indexing
        pct = s.get("percent", 0)
        eta = s.get("eta_seconds", 0)
        eta_str = f"{eta // 60}m {eta % 60}s" if eta else "calculating..."
        return (
            f"Indexing: {s['processed_files']}/{s['files_to_process']} files "
            f"({pct}%) — ETA: {eta_str}\n"
            f"Current: {s['current_file']}"
        )

    async def start(self, reset: bool = False, errors_only: bool = False) -> str:
        """Run the indexing job. Returns a summary string when done.
        errors_only: Re-index only files that had errors in the last run (fast, no full scan).
        """
        if self.status in ("preparing", "scanning", "indexing"):
            return "An indexing job is already running. Use stop_indexing() first."

        self._cancelled = False
        self.errors = []
        self.error_files = []
        self.processed_files = 0
        self.files_deleted = 0
        self.current_file = ""
        self.parse_time_seconds = 0.0
        self.embed_time_seconds = 0.0

        try:
            self.status = "preparing"
            self.detail_status = "Loading settings..."
            logging.info("Indexing job started: Preparing..." + (" (errors only)" if errors_only else ""))

            # Load settings
            user_settings = settings_module.load_settings()
            directories = user_settings["directories"]
            extensions = user_settings["extensions"]
            exclusions = user_settings["exclusions"]

            if not directories and not errors_only:
                self.status = "error"
                self.errors.append("No directories configured. Use add_directory() first.")
                return self.errors[-1]

            # Run blocking I/O in executor
            loop = asyncio.get_running_loop()

            # Quiesce WebUI + MCP so Ollama only runs embeddings (no chat/Q&A contention)
            self.detail_status = "Stopping WebUI and MCP for indexing..."
            _quiesce_services_for_indexing()

            # Attempt to unload other models to free VRAM
            try:
                self.detail_status = "Checking active models..."
                import httpx
                async with httpx.AsyncClient() as client:
                    resp = await client.get(f"{config.OLLAMA_BASE_URL}/api/ps")
                    if resp.status_code == 200:
                        data = resp.json()
                        models = data.get("models", [])
                        for m in models:
                            name = m["name"]
                            # precise check or loose check? 'nomic-embed-text' vs 'nomic-emb'
                            # If it's NOT the embedding model, unload it.
                            # config.EMBEDDING_MODEL might be 'nomic-embed-text'
                            # name might be 'nomic-embed-text:latest'
                            if config.EMBEDDING_MODEL not in name and "nomic" not in name:
                                self.detail_status = f"Unloading LLM {name}..."
                                logging.info(f"Unloading LLM {name} to free VRAM for indexing...")
                                await client.post(f"{config.OLLAMA_BASE_URL}/api/generate", 
                                                json={"model": name, "keep_alive": 0})
            except Exception as e:
                logging.warning(f"Failed to unload other models: {e}")

            # Import heavy deps
            self.detail_status = "Importing heavy libraries (LangChain, Qdrant)..."
            from langchain_text_splitters import RecursiveCharacterTextSplitter
            from langchain_qdrant import QdrantVectorStore
            from qdrant_client import QdrantClient
            from qdrant_client.models import Distance, VectorParams
            from llm_backend import get_embeddings

            # Unload LLM to prevent VRAM contention
            self.detail_status = "Unloading Chat LLM..."
            unload_llm()
            await asyncio.sleep(2)  # Allow Ollama to release VRAM

            self.detail_status = "Initializing Embeddings..."
            embeddings = get_embeddings()
            emb_url = getattr(embeddings, "base_url", None) or getattr(config, "OLLAMA_BASE_URL", "?")
            logging.info(f"Indexer embedding: base_url={emb_url}, model={config.EMBEDDING_MODEL}")
            if "11434" not in str(emb_url) and "localhost" in str(emb_url):
                logging.warning(
                    f"Embedding URL {emb_url} does not use port 11434 (Ollama default). "
                    "If you expect Ollama, check OLLAMA_BASE_URL in settings."
                )

            # Pre-index resource check: test embedding to detect insufficient RAM / overload
            self.detail_status = "Checking embedding availability..."
            try:
                _ = await loop.run_in_executor(None, lambda: embeddings.embed_query("test"))
                logging.info("Embedding pre-flight check passed.")
            except Exception as e:
                err = str(e)
                if "EOF" in err or "500" in err:
                    logging.error(
                        f"Embedding pre-flight FAILED (EOF/500): {e}. "
                        "Ollama may be overloaded, wrong URL, or another service is responding. "
                        "Check OLLAMA_BASE_URL and ensure Ollama is running on port 11434."
                    )
                logging.warning(
                    f"Insufficient resources or embedding service unavailable: {e}. "
                    "Indexing will proceed with adaptive backoff; many files may fail."
                )


            # Handle reset or integrity check
            self.detail_status = "Checking Index Integrity..."
            if reset and not errors_only:
                try:
                    client = QdrantClient(url=config.QDRANT_URL)
                    client.delete_collection(config.COLLECTION_NAME)
                    logging.warning(f"Deleted Qdrant collection: {config.COLLECTION_NAME}")
                except Exception:
                    pass  # collection may not exist yet
                if STATE_FILE.exists():
                    STATE_FILE.unlink()
                if ERROR_FILES_FILE.exists():
                    ERROR_FILES_FILE.unlink()
                    logging.info("Cleared indexing_errors.json")
                if logging_config.clear_log_file():
                    logging.info("Log file cleared for fresh re-index.")
            elif not errors_only:
                # Integrity Check: If state exists but Qdrant is missing/empty, force partial reset
                try:
                    client = QdrantClient(url=config.QDRANT_URL)
                    try:
                        info = client.get_collection(config.COLLECTION_NAME)
                        if info.points_count == 0 and STATE_FILE.exists():
                            logging.warning("Index integrity check failed: Qdrant empty but state file exists. Forcing re-index.")
                            STATE_FILE.unlink()
                    except Exception:
                        logging.warning("Index integrity check failed: Qdrant collection missing. Forcing re-index.")
                        if STATE_FILE.exists():
                            STATE_FILE.unlink()
                except Exception as e:
                    logging.error(f"Failed to perform integrity check: {e}")

            state = _load_state()
            files_to_process: List[Path] = []
            files_to_delete: List[str] = []

            # Always re-scan directories first (detects deleted/moved files)
            self.status = "scanning"
            self.detail_status = "Scanning directories..."
            scan_start = time.time()
            logging.info("Starting file scan...")
            all_files = await loop.run_in_executor(
                None, _scan_files, directories, extensions, exclusions
            )
            scan_duration = time.time() - scan_start
            logging.info(f"File scan completed in {scan_duration:.4f}s. Found {len(all_files)} files.")
            self.total_files = len(all_files)

            if self._cancelled:
                return self.get_status_text()

            files_seen = {str(fp) for fp, _ in all_files}
            files_to_delete = [f for f in state.keys() if f not in files_seen]
            if files_to_delete:
                logging.info(f"Removing {len(files_to_delete)} deleted/moved files from index.")

            if errors_only:
                # Re-index only files that had errors (and still exist)
                error_paths = _load_error_files()
                for p in error_paths:
                    path = Path(p)
                    if path.exists() and path.is_file():
                        files_to_process.append(path)
                    else:
                        logging.info(f"Skipping missing/deleted error file: {p}")
                logging.info(f"Re-indexing {len(files_to_process)} error files.")
            else:
                # Full run: process new or modified files
                for file_path, mtime in all_files:
                    str_path = str(file_path)
                    try:
                        if str_path not in state or state[str_path].get("mtime") != mtime:
                            files_to_process.append(file_path)
                    except OSError:
                        logging.warning(f"Could not access {file_path}")

            self.files_to_process = len(files_to_process)

            if not files_to_process and not files_to_delete:
                _save_state(state)
                _save_index_config_snapshot(directories, extensions, exclusions)
                self.status = "complete"
                return f"No changes detected. Index is up to date ({self.total_files} files)."

            self.status = "indexing"
            self._start_time = time.time()
            logging.info(f"Starting indexing of {self.files_to_process} files (modified or new).")

            # Initialize vector store
            # Ensure collection exists
            client = QdrantClient(url=config.QDRANT_URL)
            if not client.collection_exists(config.COLLECTION_NAME):
                # Derive embedding dimension from model at runtime
                test_vector = await loop.run_in_executor(None, lambda: embeddings.embed_query("test"))
                vector_size = len(test_vector)
                logging.info(f"Creating Qdrant collection with vector size {vector_size} (from {config.EMBEDDING_MODEL})")
                client.create_collection(
                    collection_name=config.COLLECTION_NAME,
                    vectors_config=VectorParams(size=vector_size, distance=Distance.COSINE),
                )
            vectorstore = QdrantVectorStore.from_existing_collection(
                embedding=embeddings,
                collection_name=config.COLLECTION_NAME,
                url=config.QDRANT_URL,
            )

            # Process deletions
            if files_to_delete:
                for file_path in files_to_delete:
                    chunk_ids = state[file_path].get("chunk_ids", [])
                    if chunk_ids:
                        try:
                            await loop.run_in_executor(None, vectorstore.delete, chunk_ids)
                        except Exception as e:
                            logging.error(f"Error deleting chunks for {file_path}: {e}")
                    del state[file_path]
                    self.files_deleted += 1
                _save_state(state)

            if self._cancelled:
                return self.get_status_text()

            # Process new/modified files
            text_splitter = RecursiveCharacterTextSplitter(
                chunk_size=config.CHUNK_SIZE,
                chunk_overlap=config.CHUNK_OVERLAP,
                separators=["\n\n", "\n", ". ", " ", ""],
            )

            state_save_interval = 50  # save state every N files
            files_since_save = 0
            # Adaptive backoff: increase delay, decrease batch size on 500/EOF
            # Errors-only: start conservative (these files failed before)
            base_delay = getattr(config, "INDEXING_DELAY", 0.3)
            base_batch = getattr(config, "INDEXING_BATCH_SIZE", 5)
            if errors_only:
                effective_delay = [max(base_delay * 2, 1.0)]  # At least 1s between batches
                effective_batch_size = [min(base_batch, 2)]    # Max 2 chunks at a time
                logging.info(f"Re-indexing errors with conservative settings: delay={effective_delay[0]:.1f}s, batch_size={int(effective_batch_size[0])}")
            else:
                effective_delay = [base_delay]
                effective_batch_size = [base_batch]

            for file_path in files_to_process:
                if self._cancelled:
                    break

                str_path = str(file_path)
                self.current_file = file_path.name

                try:
                    t_parse_start = time.time()
                    content = await loop.run_in_executor(None, _read_file_content, file_path)

                    # Remove old chunks if updating
                    if str_path in state:
                        old_ids = state[str_path].get("chunk_ids", [])
                        if old_ids:
                            await loop.run_in_executor(None, vectorstore.delete, old_ids)

                    if not content or len(content.strip()) == 0:
                        # No text found (image-only PDF, empty file, etc.) — not an error, just skip
                        self.parse_time_seconds += time.time() - t_parse_start
                        logging.info(f"No text in {file_path.name} (ext={file_path.suffix})")
                        state[str_path] = {"mtime": os.path.getmtime(file_path), "chunk_ids": []}
                        self.processed_files += 1
                        continue

                    max_chars = getattr(config, "INDEXING_MAX_FILE_CHARS", 500000)
                    if len(content) > max_chars:
                        content = content[:max_chars]
                        logging.info(f"Truncated {file_path.name} to {max_chars:,} chars (was larger)")

                    current_mtime = os.path.getmtime(file_path)
                    chunks = text_splitter.create_documents(
                        [content],
                        metadatas=[{
                            "source": str_path,
                            "filename": file_path.name,
                            "file_modified": time.ctime(current_mtime),
                        }],
                    )

                    self.parse_time_seconds += time.time() - t_parse_start

                    # Generate IDs for all chunks in this file
                    file_new_ids = []
                    for j, chunk in enumerate(chunks):
                        chunk.metadata["chunk_index"] = j
                        id_str = f"{str_path}_{current_mtime}_{j}"
                        chunk_id = hashlib.sha256(id_str.encode()).hexdigest()[:32]
                        file_new_ids.append(chunk_id)

                    # Batch-add chunks to Qdrant
                    # Adaptive Batching: Try desired batch size, split on failure
                    target_batch_size = max(1, int(effective_batch_size[0]))
                    chunk_limit = getattr(config, "INDEXING_MAX_BATCH_TOKENS", 8000)
                    consecutive_embed_failures = [0]  # Circuit breaker: skip file after N failures

                    def _backoff_on_failure():
                        """Increase delay and decrease batch size on 500/EOF."""
                        effective_delay[0] = min(effective_delay[0] * 1.5, 5.0)
                        effective_batch_size[0] = max(1, effective_batch_size[0] / 2)
                        logging.info(f"Adaptive backoff: delay={effective_delay[0]:.2f}s, batch_size={int(effective_batch_size[0])}")

                    def _maybe_circuit_break(orig_exc: Exception) -> None:
                        """After 3 consecutive EOF/500 failures, abort file to avoid grinding."""
                        err = str(orig_exc)
                        if "EOF" not in err and "500" not in err and "Connection refused" not in err:
                            return
                        consecutive_embed_failures[0] += 1
                        if consecutive_embed_failures[0] >= 3:
                            raise EmbeddingCircuitBreaker(
                                f"Aborting file after {consecutive_embed_failures[0]} consecutive embedding failures (EOF/500). "
                                "Skipping remaining chunks. Check Ollama and OLLAMA_BASE_URL."
                            ) from orig_exc

                    # Helper function for recursive adaptive batching
                    async def _adaptive_embed(docs_subset, ids_subset, depth=0):
                        if not docs_subset:
                            return

                        # 1. Check character limit first (simple heuristic for tokens)
                        total_chars = sum(len(d.page_content) for d in docs_subset)
                        # If over limit and we have more than 1 item, split immediately
                        if total_chars > chunk_limit and len(docs_subset) > 1:
                            mid = len(docs_subset) // 2
                            await _adaptive_embed(docs_subset[:mid], ids_subset[:mid], depth+1)
                            await _adaptive_embed(docs_subset[mid:], ids_subset[mid:], depth+1)
                            return

                        # 2. Try to embed current batch
                        try:
                            batch_delay = effective_delay[0]
                            if batch_delay > 0:
                                await asyncio.sleep(batch_delay)

                            t_embed_start = time.time()
                            await loop.run_in_executor(
                                None,
                                lambda d=docs_subset, ids=ids_subset: vectorstore.add_documents(
                                    documents=d, ids=ids
                                ),
                            )
                            self.embed_time_seconds += time.time() - t_embed_start
                            consecutive_embed_failures[0] = 0
                            return

                        except Exception as e:
                            err_msg = str(e)
                            if "EOF" in err_msg or "500" in err_msg or "Connection refused" in err_msg:
                                _backoff_on_failure()
                            # If batch size is 1, we can't split further. Retrying with backoff is the only option.
                            if len(docs_subset) == 1:
                                if "EOF" in err_msg or "500" in err_msg or "Connection refused" in err_msg:
                                    # Serious error on single item -> One retry then truncation (avoid grinding)
                                    max_retries = 1
                                    for attempt in range(max_retries):
                                        wait_time = 1 * (2 ** attempt)  # 1, 2 seconds
                                        logging.warning(f"Single item embedding failed (attempt {attempt+1}/{max_retries}), retrying in {wait_time}s: {e}")
                                        await asyncio.sleep(wait_time)
                                        try:
                                            t_embed_start = time.time()
                                            await loop.run_in_executor(
                                                None,
                                                lambda d=docs_subset, ids=ids_subset: vectorstore.add_documents(
                                                    documents=d, ids=ids
                                                ),
                                            )
                                            self.embed_time_seconds += time.time() - t_embed_start
                                            consecutive_embed_failures[0] = 0
                                            return  # Success on retry
                                        except Exception as final_e:
                                            if attempt == max_retries - 1:
                                                # Last resort: try truncated chunk (embedding API may reject long text)
                                                doc = docs_subset[0]
                                                if len(doc.page_content) > 400:
                                                    trunc_content = doc.page_content[:400]
                                                    if trunc_content:
                                                        trunc_doc = Document(page_content=trunc_content, metadata=dict(doc.metadata))
                                                        try:
                                                            await asyncio.sleep(1)  # Brief pause before truncation fallback
                                                            t_embed_start = time.time()
                                                            await loop.run_in_executor(
                                                                None,
                                                                lambda: vectorstore.add_documents(
                                                                    documents=[trunc_doc], ids=ids_subset
                                                                ),
                                                            )
                                                            self.embed_time_seconds += time.time() - t_embed_start
                                                            consecutive_embed_failures[0] = 0
                                                            logging.info(f"Embedded truncated chunk ({len(trunc_content)} chars) for {doc.metadata.get('filename', '?')}")
                                                            return
                                                        except Exception:
                                                            pass
                                                _maybe_circuit_break(final_e)
                                                raise final_e
                                else:
                                    _maybe_circuit_break(e)
                                    raise e
                            else:
                                # Batch > 1 failed -> Split and recurse (Divide and Conquer)
                                logging.warning(f"Batch of {len(docs_subset)} items failed (chars={total_chars}). Splitting and retrying...")
                                mid = len(docs_subset) // 2
                                await _adaptive_embed(docs_subset[:mid], ids_subset[:mid], depth+1)
                                await _adaptive_embed(docs_subset[mid:], ids_subset[mid:], depth+1)

                    # Execute adaptive batching for the whole file's chunks
                    # We still chunk the initial loop by target_batch_size to avoid passing 1000 items to recursion start
                    for i in range(0, len(chunks), target_batch_size):
                        batch_docs = chunks[i:i + target_batch_size]
                        batch_ids = file_new_ids[i:i + target_batch_size]
                        
                        await _adaptive_embed(batch_docs, batch_ids)

                    state[str_path] = {"mtime": current_mtime, "chunk_ids": file_new_ids}
                    files_since_save += 1

                    # Save state and errors periodically (survives restarts)
                    if files_since_save >= state_save_interval:
                        _save_state(state)
                        _save_error_files(self.error_files)
                        files_since_save = 0



                except EmbeddingCircuitBreaker as e:
                    error_msg = f"Skipped {file_path.name} (circuit breaker): {e}"
                    logging.warning(error_msg)
                    self.errors.append(error_msg)
                    if str_path not in self.error_files:
                        self.error_files.append(str_path)
                    try:
                        state[str_path] = {"mtime": os.path.getmtime(file_path), "chunk_ids": []}
                        _save_state(state)
                    except OSError:
                        pass
                    _save_error_files(self.error_files)
                except Exception as e:
                    error_msg = f"Error processing {file_path.name}: {e}"
                    logging.error(error_msg)
                    self.errors.append(error_msg)
                    if str_path not in self.error_files:
                        self.error_files.append(str_path)
                    # Add to state so failed files appear in indexed list (with 0 chunks)
                    try:
                        state[str_path] = {"mtime": os.path.getmtime(file_path), "chunk_ids": []}
                        _save_state(state)
                    except OSError:
                        pass
                    _save_error_files(self.error_files)  # Persist immediately so errors survive restarts

                self.processed_files += 1

            if self._cancelled:
                _save_state(state)  # save progress before exit
                _save_error_files(self.error_files)
                if self.parse_time_seconds > 0 or self.embed_time_seconds > 0:
                    logging.info(f"Indexing cancelled. Time breakdown: parse {self.parse_time_seconds:.0f}s, embed {self.embed_time_seconds:.0f}s")
                _restore_services_after_indexing()  # Restore services even on cancel
                return self.get_status_text()

            _save_state(state)  # final save
            _save_error_files(self.error_files)
            _save_index_config_snapshot(directories, extensions, exclusions)
            self.status = "complete"
            if self.parse_time_seconds > 0 or self.embed_time_seconds > 0:
                logging.info(f"Indexing complete. Time breakdown: parse {self.parse_time_seconds:.0f}s, embed {self.embed_time_seconds:.0f}s")
            _restore_services_after_indexing()  # Warm LLM, restart WebUI + MCP
            return self.get_status_text()

        except Exception as e:
            self.status = "error"
            error_msg = f"Indexing failed: {e}"
            self.errors.append(error_msg)
            _save_error_files(self.error_files)
            logging.error(error_msg)
            return error_msg


# Singleton job instance for the MCP server
_current_job: Optional[IndexingJob] = None


async def _index_single_file_async(file_path: Path) -> Tuple[bool, str]:
    """Index a single file. Does not quiesce services. Returns (success, message)."""
    job = get_current_job()
    if job.status in ("preparing", "scanning", "indexing"):
        return False, "An indexing job is already running. Wait for it to finish."

    str_path = str(file_path)
    if not file_path.exists() or not file_path.is_file():
        return False, f"File not found: {str_path}"

    user_settings = settings_module.load_settings()
    extensions = user_settings.get("extensions") or user_settings.get("SUPPORTED_EXTENSIONS") or []
    if file_path.suffix.lower() not in [e.lower() for e in extensions]:
        return False, f"File type {file_path.suffix} is not in supported extensions."

    loop = asyncio.get_running_loop()
    from langchain_qdrant import QdrantVectorStore
    from qdrant_client import QdrantClient
    from qdrant_client.models import VectorParams, Distance
    from langchain_text_splitters import RecursiveCharacterTextSplitter
    from llm_backend import get_embeddings

    embeddings = get_embeddings()
    client = QdrantClient(url=config.QDRANT_URL)
    if not client.collection_exists(config.COLLECTION_NAME):
        test_vector = embeddings.embed_query("test")
        client.create_collection(
            collection_name=config.COLLECTION_NAME,
            vectors_config=VectorParams(size=len(test_vector), distance=Distance.COSINE),
        )
    vectorstore = QdrantVectorStore.from_existing_collection(
        embedding=embeddings,
        collection_name=config.COLLECTION_NAME,
        url=config.QDRANT_URL,
    )

    content = await loop.run_in_executor(None, _read_file_content, file_path)
    state = _load_state()

    # Delete old chunks
    if str_path in state:
        old_ids = state[str_path].get("chunk_ids", [])
        if old_ids:
            try:
                await loop.run_in_executor(None, vectorstore.delete, old_ids)
            except Exception as e:
                logging.warning(f"Could not delete old chunks for {str_path}: {e}")

    if not content or len(content.strip()) == 0:
        current_mtime = os.path.getmtime(file_path)
        state[str_path] = {"mtime": current_mtime, "chunk_ids": []}
        _save_state(state)
        error_files = _load_error_files()
        if str_path in error_files:
            error_files.remove(str_path)
            _save_error_files(error_files)
        return True, "File has no extractable text. Chunks cleared from index."

    max_chars = getattr(config, "INDEXING_MAX_FILE_CHARS", 500000)
    if len(content) > max_chars:
        content = content[:max_chars]
        logging.info(f"Truncated {file_path.name} to {max_chars:,} chars (single-file re-index)")

    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=config.CHUNK_SIZE,
        chunk_overlap=config.CHUNK_OVERLAP,
        separators=["\n\n", "\n", ". ", " ", ""],
    )
    current_mtime = os.path.getmtime(file_path)
    chunks = text_splitter.create_documents(
        [content],
        metadatas=[{"source": str_path, "filename": file_path.name, "file_modified": time.ctime(current_mtime)}],
    )
    file_new_ids = []
    for j, chunk in enumerate(chunks):
        chunk.metadata["chunk_index"] = j
        id_str = f"{str_path}_{current_mtime}_{j}"
        chunk_id = hashlib.sha256(id_str.encode()).hexdigest()[:32]
        file_new_ids.append(chunk_id)

    # Add in batches of 5 to avoid overload
    batch_size = 5
    for i in range(0, len(chunks), batch_size):
        batch_docs = chunks[i : i + batch_size]
        batch_ids = file_new_ids[i : i + batch_size]
        await loop.run_in_executor(
            None,
            lambda d=batch_docs, ids=batch_ids: vectorstore.add_documents(documents=d, ids=ids),
        )
        if i + batch_size < len(chunks):
            await asyncio.sleep(0.5)

    state[str_path] = {"mtime": current_mtime, "chunk_ids": file_new_ids}
    _save_state(state)
    error_files = _load_error_files()
    if str_path in error_files:
        error_files.remove(str_path)
        _save_error_files(error_files)
    return True, f"Re-indexed {file_path.name}: {len(chunks)} chunks."


def index_single_file(path: str) -> Tuple[bool, str]:
    """Synchronous wrapper to index a single file. Runs in a new event loop."""
    try:
        return asyncio.run(_index_single_file_async(Path(path)))
    except Exception as e:
        logging.error(f"Single-file index failed for {path}: {e}")
        return False, str(e)


def get_current_job() -> IndexingJob:
    """Get or create the singleton IndexingJob."""
    global _current_job
    if _current_job is None:
        _current_job = IndexingJob()
    return _current_job
